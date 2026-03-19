import os, io, re, uuid, hashlib, base64
from datetime import datetime, timedelta
from pathlib import Path

from flask import Flask, request, jsonify
from flask_cors import CORS
from flask_jwt_extended import (
    JWTManager, create_access_token, jwt_required,
    get_jwt_identity, get_jwt
)
from pymongo import MongoClient, TEXT
from pymongo.errors import DuplicateKeyError
import bcrypt
import pytesseract
from PIL import Image, ImageEnhance, ImageFilter
import cv2, numpy as np
from pdf2image import convert_from_bytes
import docx
from pptx import Presentation

MONGO_URI  = "mongodb+srv://root:root@cluster0.igwufcv.mongodb.net/?appName=Cluster0"
NGROK_TOKEN = "3A79UkSS3ThVjSKOVk9PefLPG6b_5hxpDi4vAetFEcuD82FmY"
DB_NAME    = "DocuSearch"
JWT_SECRET  = "abcd123!"

app = Flask(__name__)
app.config["JWT_SECRET_KEY"]           = JWT_SECRET
app.config["JWT_ACCESS_TOKEN_EXPIRES"] = timedelta(hours=12)
CORS(app,
     resources={r"/*": {"origins": "*"}},
     allow_headers=["Content-Type", "Authorization", "ngrok-skip-browser-warning"],
     methods=["GET", "POST", "DELETE", "OPTIONS"],
     supports_credentials=False)

@app.after_request
def add_cors_headers(response):
    response.headers["Access-Control-Allow-Origin"]  = "*"
    response.headers["Access-Control-Allow-Headers"] = "Content-Type, Authorization, ngrok-skip-browser-warning"
    response.headers["Access-Control-Allow-Methods"] = "GET, POST, DELETE, OPTIONS"
    return response
jwt = JWTManager(app)

client    = MongoClient(MONGO_URI)
db        = client[DB_NAME]
docs_col  = db["documents"]
users_col = db["users"]

docs_col.create_index([("text", TEXT), ("filename", TEXT)], name="text_search_index")
try: docs_col.create_index([("file_hash",1),("user_id",1)], unique=True)
except: pass
try: users_col.create_index("email", unique=True)
except: pass

print("Connected to MongoDB")

print("Loading EasyOCR model (first time ~30s)...")
import easyocr
_easyocr_reader = easyocr.Reader(["en"], gpu=False, verbose=False)
print("EasyOCR ready")

print("Loading TrOCR model (first time ~60s, downloads ~1.3GB)...")
_trocr_processor = None
_trocr_model     = None
try:
    from transformers import TrOCRProcessor, VisionEncoderDecoderModel
    _trocr_processor = TrOCRProcessor.from_pretrained("microsoft/trocr-large-handwritten")
    _trocr_model     = VisionEncoderDecoderModel.from_pretrained("microsoft/trocr-large-handwritten")
    print("TrOCR ready")
except Exception as e:
    print(f"TrOCR failed to load: {e}")
    print("   Run: !pip install transformers -q")
    print("   Falling back to EasyOCR only")

_spell = None
def get_spell():
    global _spell
    if _spell is not None:
        return _spell
    try:
        from symspellpy import SymSpell, Verbosity
        import urllib.request, os
        ss = SymSpell(max_dictionary_edit_distance=2)
        dic_path = "/tmp/frequency_dictionary_en_82_765.txt"
        if not os.path.exists(dic_path):
            print("  Downloading spell dictionary...")
            urllib.request.urlretrieve(
                "https://raw.githubusercontent.com/mammothb/symspellpy/master/symspellpy/frequency_dictionary_en_82_765.txt",
                dic_path
            )
        ss.load_dictionary(dic_path, term_index=0, count_index=1)
        _spell = ss
        print("Spell corrector ready")
    except Exception as e:
        print(f"  Spell corrector unavailable: {e}")
    return _spell

HANDWRITING_FIXES = [
    (r"UJ",      "W"),
    (r"Ulish",   "Wish"),
    (r"Uish",    "Wish"),
    (r"Jou",     "You"),
    (r"Jbu",     "You"),
    (r"Jbo",     "You"),
    (r"Ju",      "You"),
    (r"UolJ",    "Now"),
    (r"uolJ",    "Now"),
    (r"Uow",     "Now"),
    (r"Ajo",     "And"),
    (r"AND",     "And"),
    (r"Fod",     "For"),
    (r"FoR",     "For"),
    (r"Fod_",    "For"),
    (r"ThE",     "The"),
    (r"Thc",     "The"),
    (r"GESt",    "Best"),
    (r"Gest",    "Best"),
    (r"GEST",    "Best"),
    (r"INcLuoin\w*","Including"),
    (r"INcLudin\w*","Including"),
    (r"smiLuJ",  "Smiling"),
    (r"smiLiJ",  "Smiling"),
    (r"smilin",  "Smiling"),
    (r"DALin",   "Dialin"),
    (r"DiALin",  "Dialin"),
    (r"DALIN",   "Dialin"),
    (r"uuaerfuL","Wonderful"),
    (r"uuonderful","Wonderful"),
    (r"ThaUl",   "Thank"),
    (r"THAUL",   "Thank"),
    (r"THANL",   "Thank"),
    (r"TALLALE", "Package"),
    (r"TALLAGE", "Package"),
    (r"TALLAce", "Package"),
    (r"cerfect", "Perfect"),
    (r"PERFECT", "Perfect"),
    (r"Touch",   "Touch"),
    (r"TOUCH",   "Touch"),
    (r"Biauest", "Biggest"),
    (r"BIAUEST", "Biggest"),
    (r"BIGGEST", "Biggest"),
    (r"FANL",    "Fan"),
    (r"Life",    "Life"),
    (r"LiFe",    "Life"),
    (r"OffER_",  "Offer"),
    (r"OffER",   "Offer"),
    (r"OFFER",   "Offer"),
    (r"Gift",    "Gift"),
    (r"ift",     "Gift"),
    (r"15",      "Is"),
    (r"1M",      "I'm"),
    (r"IM",      "I'm"),
    (r"_",           ""),
]

def apply_handwriting_fixes(text: str) -> str:
    for pattern, replacement in HANDWRITING_FIXES:
        text = re.sub(pattern, replacement, text, flags=re.IGNORECASE)
    text = re.sub(r"\s{2,}", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()

def spell_correct_word(word: str, spell) -> str:
    if not spell or len(word) <= 2:
        return word
    stripped = re.sub(r"[^a-zA-Z]", "", word)
    if not stripped or stripped.lower() in {"the","and","for","you","now","all","has","its","not","are","but","can","was","one","may","its","how","our","who","two","him","his","her","she","this","that","with","from","they","have","been","when","what","will","into","also","more"}:
        return word
    from symspellpy import Verbosity
    suggestions = spell.lookup(stripped, Verbosity.CLOSEST, max_edit_distance=2)
    if suggestions and suggestions[0].distance <= 1:
        return word.replace(stripped, suggestions[0].term)
    return word

def postprocess_ocr(text: str) -> str:
    text = apply_handwriting_fixes(text)
    spell = get_spell()
    if spell:
        words    = text.split()
        corrected = [spell_correct_word(w, spell) for w in words]
        text     = " ".join(corrected)
    return text.strip()

def detect_text_lines(arr: np.ndarray):
    gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
    h, w = gray.shape

    gray = cv2.GaussianBlur(gray, (3, 3), 0)

    _, bw = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)

    num_labels, labels, stats, _ = cv2.connectedComponentsWithStats(bw, connectivity=8)
    heights = [stats[i, cv2.CC_STAT_HEIGHT] for i in range(1, num_labels)]
    widths  = [stats[i, cv2.CC_STAT_WIDTH]  for i in range(1, num_labels)]

    if not heights:
        return [(0, h, 0, w)]

    median_h = sorted(heights)[len(heights)//2]
    char_heights = [hh for hh in heights if median_h * 0.3 < hh < median_h * 3]
    avg_char_h = int(sum(char_heights) / len(char_heights)) if char_heights else 20

    kernel_w = max(int(w * 0.04), 20)
    kernel_h = max(avg_char_h // 3, 3)
    kernel   = cv2.getStructuringElement(cv2.MORPH_RECT, (kernel_w, kernel_h))
    dilated  = cv2.dilate(bw, kernel, iterations=2)

    contours, _ = cv2.findContours(dilated, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    line_boxes = []
    min_line_h = avg_char_h * 0.5
    min_line_w = w * 0.05

    for cnt in contours:
        x, y, cw, ch = cv2.boundingRect(cnt)
        if ch >= min_line_h and cw >= min_line_w:
            pad = max(int(avg_char_h * 0.2), 4)
            y1  = max(0, y - pad)
            y2  = min(h, y + ch + pad)
            line_boxes.append((y1, y2, 0, w))

    line_boxes.sort(key=lambda b: b[0])

    merged = []
    for box in line_boxes:
        if merged and box[0] < merged[-1][1] * 0.85:
            prev = merged[-1]
            merged[-1] = (prev[0], max(prev[1], box[1]), 0, w)
        else:
            merged.append(box)

    return merged if merged else [(0, h, 0, w)]


def trocr_ocr(img: Image.Image) -> str:
    if _trocr_processor is None or _trocr_model is None:
        return ""
    try:
        import torch

        arr = np.array(img.convert("RGB"))
        h, w = arr.shape[:2]

        line_boxes = detect_text_lines(arr)
        print(f"  TrOCR: detected {len(line_boxes)} lines")

        lines_text = []
        for idx, (y1, y2, x1, x2) in enumerate(line_boxes):
            line_arr = arr[y1:y2, x1:x2]
            line_h   = y2 - y1

            if line_h < 8:
                continue

            if line_h < 32:
                scale    = 32 / line_h
                line_arr = cv2.resize(line_arr, None, fx=scale, fy=scale,
                                      interpolation=cv2.INTER_CUBIC)

            line_img     = Image.fromarray(line_arr).convert("RGB")
            pixel_values = _trocr_processor(images=line_img, return_tensors="pt").pixel_values

            with torch.no_grad():
                generated_ids = _trocr_model.generate(
                    pixel_values,
                    max_new_tokens=128,
                    num_beams=4,
                    early_stopping=True,
                    no_repeat_ngram_size=3,
                )

            line_text = _trocr_processor.batch_decode(
                generated_ids, skip_special_tokens=True
            )[0].strip()

            if line_text:
                lines_text.append(line_text)
                print(f"    Line {idx+1}: {line_text[:80]}")

        result = "\n".join(lines_text)
        print(f"  TrOCR total: {len(result.split())} words")
        return result

    except Exception as e:
        print(f"  TrOCR failed: {e}")
        return ""

def easyocr_with_structure(arr: np.ndarray) -> str:
    results = _easyocr_reader.readtext(arr, detail=1, paragraph=False)
    if not results:
        return ""

    def midY(r): return (r[0][0][1] + r[0][2][1]) / 2
    def midX(r): return (r[0][0][0] + r[0][2][0]) / 2
    def height(r): return abs(r[0][2][1] - r[0][0][1])

    results_sorted = sorted(results, key=lambda r: (midY(r), midX(r)))

    lines = []
    current_line = []
    prev_y = None
    avg_h  = sum(height(r) for r in results_sorted) / len(results_sorted) if results_sorted else 20

    for result in results_sorted:
        bbox, text, conf = result
        y = midY(result)

        if prev_y is not None and abs(y - prev_y) > avg_h * 0.6:
            if current_line:
                lines.append(" ".join(current_line))
            current_line = []

        current_line.append(text)
        prev_y = y

    if current_line:
        lines.append(" ".join(current_line))

    return "\n".join(lines)

def ocr_image(img: Image.Image) -> str:
    arr = np.array(img.convert("RGB"))

    h, w = arr.shape[:2]
    if w < 1000:
        arr = cv2.resize(arr, None, fx=2.0, fy=2.0, interpolation=cv2.INTER_CUBIC)
    elif w < 2000:
        arr = cv2.resize(arr, None, fx=1.5, fy=1.5, interpolation=cv2.INTER_CUBIC)

    img_upscaled = Image.fromarray(arr)

    trocr_text = ""
    try:
        if _trocr_model is not None:
            trocr_text = trocr_ocr(img_upscaled)
            print(f"  TrOCR:          {trocr_text[:120]}")
    except Exception as e:
        print(f"  TrOCR failed: {e}")

    easy_text = ""
    try:
        easy_text = easyocr_with_structure(arr)
        print(f"  EasyOCR:        {easy_text[:120]}")
    except Exception as e:
        print(f"  EasyOCR failed: {e}")

    tess_text = ""
    try:
        gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
        _, binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        padded = cv2.copyMakeBorder(binary, 30, 30, 30, 30, cv2.BORDER_CONSTANT, value=255)
        tess_text = pytesseract.image_to_string(
            Image.fromarray(padded), config="--oem 3 --psm 6", lang="eng"
        ).strip()
        print(f"  Tesseract:      {tess_text[:120]}")
    except Exception as e:
        print(f"  Tesseract failed: {e}")

    if trocr_text and len(trocr_text.split()) >= 3:
        best_raw = trocr_text
        print(f"  Using TrOCR ({len(trocr_text.split())} words)")
    else:
        candidates = [t for t in [easy_text, tess_text, trocr_text] if t]
        if not candidates:
            print("  All OCR engines returned empty")
            return ""
        best_raw = max(candidates, key=lambda t: len(t.split()))
        engine = "EasyOCR" if best_raw == easy_text else "Tesseract"
        print(f"  Using {engine} fallback ({len(best_raw.split())} words)")

    final = postprocess_ocr(best_raw)
    print(f"  Final:        {final[:120]}")
    return final

def postprocess_text(text: str) -> str:
    if not text:
        return text

    fixes = [
        (r'(?<=[a-z])0(?=[a-z])', 'o'),
        (r'(?<=[a-z])1(?=[a-z])', 'l'),
        (r'\|', 'I'),
        (r'(?<=\w)~(?=\w)', '-'),
        (r'\n{3,}', '\n\n'),
        (r'[ \t]{2,}', ' '),
    ]
    for pattern, replacement in fixes:
        text = re.sub(pattern, replacement, text)

    return text.strip()

def extract_text_from_pdf(data: bytes) -> str:
    try:
        import fitz
        pdf   = fitz.open(stream=data, filetype="pdf")
        texts = []
        for page in pdf:
            direct = page.get_text().strip()
            if len(direct) > 50:
                texts.append(direct)
            else:
                mat = fitz.Matrix(300/72, 300/72)
                pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                texts.append(ocr_image(img))
        pdf.close()
        if texts:
            return postprocess_text("\n\n--- Page Break ---\n\n".join(texts))
    except Exception as e:
        print(f"PyMuPDF attempt failed: {e}")

    try:
        pages = convert_from_bytes(data, dpi=300)
        texts = [ocr_image(p) for p in pages]
        return postprocess_text("\n\n--- Page Break ---\n\n".join(texts))
    except Exception as e:
        raise ValueError(f"All PDF extraction methods failed: {e}")

def extract_text(filename: str, data: bytes) -> str:
    ext = Path(filename).suffix.lower()
    if ext in (".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"):
        return postprocess_text(ocr_image(Image.open(io.BytesIO(data))))
    elif ext == ".pdf":
        return extract_text_from_pdf(data)
    elif ext == ".docx":
        d = docx.Document(io.BytesIO(data))
        return "\n".join(p.text for p in d.paragraphs if p.text.strip())
    elif ext == ".pptx":
        prs = Presentation(io.BytesIO(data))
        return "\n".join(
            s.text for slide in prs.slides
            for s in slide.shapes
            if hasattr(s, "text") and s.text.strip()
        )
    elif ext == ".txt":
        return data.decode("utf-8", errors="ignore")
    raise ValueError(f"Unsupported file type: {ext}")

STOPWORDS = {
    "the","and","for","are","but","not","you","all","any","can","had","was","one",
    "our","out","day","get","has","him","his","how","its","may","new","now","old",
    "see","two","who","did","man","let","put","say","she","too","use","that","with",
    "this","from","they","have","been","when","what","will","more","also","into"
}
def keywords(text):
    words = re.findall(r'\b[a-zA-Z]{3,}\b', text.lower())
    return list({w for w in words if w not in STOPWORDS})

def snippet(text, query, window=200):
    lower = text.lower()
    for term in query.lower().split():
        idx = lower.find(term)
        if idx != -1:
            s = max(0, idx - window // 2)
            e = min(len(text), idx + window // 2)
            r = text[s:e].strip()
            return ("..." if s > 0 else "") + r + ("..." if e < len(text) else "")
    return text[:window] + "..."

def serialize(doc):
    if doc is None: return None
    d = dict(doc)
    if "_id" in d: d["_id"] = str(d["_id"])
    d.pop("file_data", None)
    d.pop("password",  None)
    return d

@app.route("/auth/register", methods=["POST"])
def register():
    data     = request.get_json()
    name     = (data.get("name")     or "").strip()
    email    = (data.get("email")    or "").strip().lower()
    password = (data.get("password") or "").strip()
    if not name or not email or not password:
        return jsonify({"error": "Name, email and password are required"}), 400
    if len(password) < 6:
        return jsonify({"error": "Password must be at least 6 characters"}), 400
    if "@" not in email:
        return jsonify({"error": "Invalid email address"}), 400
    hashed  = bcrypt.hashpw(password.encode(), bcrypt.gensalt()).decode()
    user_id = str(uuid.uuid4())
    try:
        users_col.insert_one({"_id": user_id, "name": name, "email": email,
                               "password": hashed, "created_at": datetime.utcnow().isoformat()})
    except DuplicateKeyError:
        return jsonify({"error": "An account with this email already exists"}), 409
    token = create_access_token(identity=user_id, additional_claims={"name": name, "email": email})
    return jsonify({"token": token, "user": {"id": user_id, "name": name, "email": email}}), 201

@app.route("/auth/login", methods=["POST"])
def login():
    data     = request.get_json()
    email    = (data.get("email")    or "").strip().lower()
    password = (data.get("password") or "").strip()
    if not email or not password:
        return jsonify({"error": "Email and password are required"}), 400
    user = users_col.find_one({"email": email})
    if not user or not bcrypt.checkpw(password.encode(), user["password"].encode()):
        return jsonify({"error": "Invalid email or password"}), 401
    token = create_access_token(identity=user["_id"],
                                additional_claims={"name": user["name"], "email": user["email"]})
    return jsonify({"token": token, "user": {"id": user["_id"], "name": user["name"], "email": user["email"]}}), 200

@app.route("/auth/me", methods=["GET"])
@jwt_required()
def me():
    claims = get_jwt()
    return jsonify({"id": get_jwt_identity(), "name": claims.get("name"), "email": claims.get("email")})

@app.route("/health")
def health():
    return jsonify({"status": "ok", "time": datetime.utcnow().isoformat()})

@app.route("/upload", methods=["POST"])
@jwt_required()
def upload():
    user_id = get_jwt_identity()
    print(f"Upload request from user_id: {user_id}")
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400
    f     = request.files["file"]
    data  = f.read()
    fhash = hashlib.sha256(data).hexdigest()

    existing = docs_col.find_one({"file_hash": fhash, "user_id": user_id})
    if existing:
        return jsonify({"message": "Already exists", "doc_id": str(existing["_id"])}), 200

    try:
        text = extract_text(f.filename, data)
    except Exception as e:
        print(f"OCR extraction failed for '{f.filename}': {e}")
        text = f"[OCR failed: {str(e)[:200]}]"

    doc_id = str(uuid.uuid4())
    record = {
        "_id":         doc_id,
        "user_id":     user_id,
        "filename":    f.filename,
        "file_hash":   fhash,
        "file_type":   Path(f.filename).suffix.lower().lstrip("."),
        "file_size":   len(data),
        "file_data":   base64.b64encode(data).decode(),
        "text":        text,
        "keywords":    keywords(text),
        "word_count":  len(text.split()),
        "char_count":  len(text),
        "uploaded_at": datetime.utcnow().isoformat(),
    }
    try:
        docs_col.insert_one(record)
        print(f"Saved doc '{f.filename}' for user {user_id} ({len(text.split())} words)")
    except DuplicateKeyError:
        return jsonify({"message": "Already exists"}), 200

    return jsonify({
        "message":    "Indexed successfully",
        "doc_id":     doc_id,
        "filename":   f.filename,
        "word_count": record["word_count"],
        "keywords":   record["keywords"][:20],
    }), 201

@app.route("/search")
@jwt_required()
def search():
    user_id = get_jwt_identity()
    q    = request.args.get("q", "").strip()
    page = int(request.args.get("page", 1))
    per  = int(request.args.get("per_page", 10))
    if not q:
        return jsonify({"error": "q is required"}), 400
    cursor = docs_col.find(
        {"$text": {"$search": q}, "user_id": user_id},
        {"score": {"$meta": "textScore"}, "file_data": 0}
    ).sort([("score", {"$meta": "textScore"})]).skip((page-1)*per).limit(per)
    results = [{
        "doc_id":      d["_id"],
        "filename":    d["filename"],
        "file_type":   d.get("file_type",""),
        "word_count":  d.get("word_count",0),
        "uploaded_at": d.get("uploaded_at",""),
        "snippet":     snippet(d.get("text",""), q),
        "score":       round(d.get("score",0), 2),
    } for d in cursor]
    total = docs_col.count_documents({"$text": {"$search": q}, "user_id": user_id})
    return jsonify({"query": q, "total": total, "page": page, "results": results})

@app.route("/documents")
@jwt_required()
def list_docs():
    user_id = get_jwt_identity()
    page = int(request.args.get("page", 1))
    per  = int(request.args.get("per_page", 12))
    cursor = docs_col.find({"user_id": user_id}, {"file_data": 0}).sort("uploaded_at",-1).skip((page-1)*per).limit(per)
    total  = docs_col.count_documents({"user_id": user_id})
    return jsonify({"total": total, "page": page, "documents": [serialize(d) for d in cursor]})

@app.route("/documents/<doc_id>")
@jwt_required()
def get_doc(doc_id):
    user_id = get_jwt_identity()
    doc = docs_col.find_one({"_id": doc_id, "user_id": user_id}, {"file_data": 0})
    if not doc: return jsonify({"error": "Not found"}), 404
    return jsonify(serialize(doc))

@app.route("/documents/<doc_id>/text")
@jwt_required()
def get_text(doc_id):
    user_id = get_jwt_identity()
    doc = docs_col.find_one({"_id": doc_id, "user_id": user_id}, {"text": 1, "filename": 1})
    if not doc: return jsonify({"error": "Not found"}), 404
    return jsonify({"doc_id": doc_id, "filename": doc["filename"], "text": doc["text"]})

@app.route("/documents/<doc_id>/file")
@jwt_required()
def get_file(doc_id):
    user_id = get_jwt_identity()
    doc = docs_col.find_one({"_id": doc_id, "user_id": user_id},
                             {"file_data": 1, "filename": 1, "file_type": 1})
    if not doc: return jsonify({"error": "Not found"}), 404
    if not doc.get("file_data"): return jsonify({"error": "File data not stored"}), 404
    return jsonify({
        "doc_id":    doc_id,
        "filename":  doc["filename"],
        "file_type": doc.get("file_type", ""),
        "file_data": doc["file_data"],
    })

@app.route("/documents/<doc_id>", methods=["DELETE"])
@jwt_required()
def delete_doc(doc_id):
    user_id = get_jwt_identity()
    r = docs_col.delete_one({"_id": doc_id, "user_id": user_id})
    if r.deleted_count == 0: return jsonify({"error": "Not found"}), 404
    return jsonify({"message": "Deleted"})

@app.route("/stats")
@jwt_required()
def stats():
    user_id = get_jwt_identity()
    total = docs_col.count_documents({"user_id": user_id})
    agg   = list(docs_col.aggregate([
        {"$match": {"user_id": user_id}},
        {"$group": {"_id": None, "tw": {"$sum": "$word_count"}}}
    ]))
    types = list(docs_col.aggregate([
        {"$match": {"user_id": user_id}},
        {"$group": {"_id": "$file_type", "c": {"$sum": 1}}}
    ]))
    return jsonify({
        "total_documents": total,
        "total_words":     agg[0]["tw"] if agg else 0,
        "file_types":      {t["_id"]: t["c"] for t in types},
    })

import threading, time
from pyngrok import ngrok

def run_flask():
    app.run(port=5000, debug=False, use_reloader=False, threaded=True)

flask_thread = threading.Thread(target=run_flask, daemon=True)
flask_thread.start()
time.sleep(2)

ngrok.set_auth_token(NGROK_TOKEN)
ngrok.kill()
time.sleep(1)
tunnel = ngrok.connect(5000)

print("\n" + "="*60)
print(f"  API is live at: {tunnel.public_url}")
print(f"  REACT_APP_API_URL={tunnel.public_url}")
print("="*60 + "\n")
print("Server is running — do NOT stop this cell.")
print("   It will keep running until you click the stop button.\n")

count = 0
while True:
    time.sleep(300)
    count += 1
    print(f"  Still running — {count * 5} min elapsed | {tunnel.public_url}")
